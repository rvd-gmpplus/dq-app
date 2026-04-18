#!/usr/bin/env python3
"""Extract use case text from USE-CASES/*.pptx into src/seed/useCases.draft.json.

Runs locally only. The pptx files stay gitignored; the draft JSON is a working
artefact used when curating typed seed data in src/seed/useCases.ts.

Usage:
    python scripts/extract_use_cases.py
"""
from __future__ import annotations

import json
import pathlib
import sys

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    print("Install python-pptx first: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

SRC = pathlib.Path("USE-CASES")
OUT = pathlib.Path("src/seed/useCases.draft.json")


def shape_text(shape) -> list[str]:
    """Collect non-empty text lines from a single shape, including tables."""
    lines: list[str] = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            text = "".join(run.text for run in para.runs).strip()
            if text:
                lines.append(text)
    if shape.has_table:
        for row in shape.table.rows:
            cells = []
            for cell in row.cells:
                cell_text = "\n".join(
                    "".join(run.text for run in para.runs).strip()
                    for para in cell.text_frame.paragraphs
                ).strip()
                cells.append(cell_text)
            if any(cells):
                lines.append(" | ".join(cells))
    return lines


def slide_payload(index: int, slide) -> dict:
    lines: list[str] = []
    title: str | None = None
    for shape in slide.shapes:
        text_lines = shape_text(shape)
        # Heuristic: first text shape that looks like a title
        if title is None and text_lines and shape.has_text_frame and shape.is_placeholder:
            ph = shape.placeholder_format
            if ph is not None and ph.idx == 0:
                title = text_lines[0]
        lines.extend(text_lines)

    notes = ""
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()

    return {
        "slide_index": index,
        "title": title,
        "lines": lines,
        "notes": notes,
    }


def extract(path: pathlib.Path) -> dict:
    prs = Presentation(str(path))
    slides = [slide_payload(i, s) for i, s in enumerate(prs.slides, start=1)]
    return {
        "source": path.name,
        "slide_count": len(slides),
        "slides": slides,
    }


def main() -> int:
    if not SRC.exists():
        print(f"No {SRC} folder found", file=sys.stderr)
        return 1
    OUT.parent.mkdir(parents=True, exist_ok=True)
    files = sorted(SRC.glob("*.pptx"))
    if not files:
        print(f"No .pptx files in {SRC}", file=sys.stderr)
        return 1
    payload = [extract(p) for p in files]
    OUT.write_text(
        json.dumps(payload, indent=2, ensure_ascii=False),
        encoding="utf-8",
    )
    total_slides = sum(f["slide_count"] for f in payload)
    print(f"Wrote {total_slides} slides from {len(files)} file(s) to {OUT}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
